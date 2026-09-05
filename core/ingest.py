"""Pipeline pemrosesan file yang diunggah.

Alurnya selalu lima tahap dan urutannya tidak berubah:

1. File diterima, disimpan sementara ke folder temp
2. Validasi tipe dan ukuran
3. Parsing menjadi teks
4. Metadata dan teks disimpan ke SQLite
5. File asli dihapus dari disk

File asli tidak pernah disimpan permanen — yang bertahan hanya hasil
pemrosesannya di database. Modul ini sengaja tidak mengimpor Streamlit;
progres dilaporkan lewat callback ``progress`` supaya bisa dipakai dari
UI apa pun (atau dari test).
"""

from __future__ import annotations

import re
import tempfile
from dataclasses import dataclass
from pathlib import Path
from typing import Callable

from core import db, errors, rag, vision
from core.config import get_settings

# --------------------------------------------------------------------------- #
# Batasan
# --------------------------------------------------------------------------- #
MAX_FILE_SIZE_MB = 5
MAX_FILE_SIZE = MAX_FILE_SIZE_MB * 1024 * 1024

# Batas per sekali unggah, bukan kuota permanen. Setelah satu batch selesai
# diproses, pengguna boleh mengunggah batch berikutnya tanpa perlu menghapus
# dokumen lama.
MAX_FILES_PER_BATCH = 5

IMAGE_TYPES = ("png", "jpg", "jpeg")
DOCUMENT_TYPES = ("pdf", "docx", "xlsx", "pptx")
ALLOWED_TYPES = IMAGE_TYPES + DOCUMENT_TYPES

TYPE_LABELS = {
    "png": "PNG",
    "jpg": "JPG",
    "jpeg": "JPEG",
    "pdf": "PDF",
    "docx": "DOCX",
    "xlsx": "XLSX",
    "pptx": "PPTX",
}

# Status yang mungkin tersimpan di kolom ``documents.status``
STATUS_INDEXED = "indexed"             # teks diekstrak dan chunk-nya masuk Qdrant
STATUS_PROCESSED = "processed"          # teks diekstrak, tapi belum terindeks
STATUS_NO_TEXT = "no_text"             # terbaca, tapi tidak ada teks di dalamnya
STATUS_PENDING_EXTRACTION = "pending"  # gambar yang gagal dibaca model vision
STATUS_FAILED = "failed"               # file tidak bisa dibaca


@dataclass(frozen=True)
class IngestResult:
    """Hasil satu kali pemrosesan file."""

    filename: str
    ok: bool
    message: str
    status: str = ""
    char_count: int = 0
    document_id: int | None = None
    chunk_count: int = 0


class RejectedFile(Exception):
    """File ditolak karena melanggar batasan (bukan error tak terduga)."""


class ParseFailed(Exception):
    """File lolos validasi tapi isinya tidak bisa dibaca."""


# --------------------------------------------------------------------------- #
# Helper
# --------------------------------------------------------------------------- #
def thousands(value: int) -> str:
    """Angka dengan pemisah ribuan gaya Indonesia, mis. ``12.480``."""
    return f"{value:,}".replace(",", ".")


def human_size(num_bytes: int) -> str:
    """Ukuran file dalam satuan yang enak dibaca, mis. ``1,4 MB``."""
    if num_bytes < 1024:
        return f"{num_bytes} B"
    if num_bytes < 1024 * 1024:
        return f"{round(num_bytes / 1024)} KB"
    megabytes = f"{num_bytes / (1024 * 1024):.1f}".replace(".", ",")
    return f"{megabytes} MB"


def file_extension(filename: str) -> str:
    return Path(filename).suffix.lower().lstrip(".")


def type_label(filetype: str) -> str:
    return TYPE_LABELS.get(filetype, filetype.upper())


def allowed_types_sentence() -> str:
    return ", ".join(TYPE_LABELS[t] for t in ALLOWED_TYPES)


def check_batch(count: int) -> str | None:
    """Periksa jumlah file dalam satu kali unggah.

    Mengembalikan pesan penolakan kalau melebihi jatah, atau ``None`` kalau
    lolos. Kelebihan jumlah menolak seluruh batch, tidak memproses sebagian,
    supaya pengguna sendiri yang memilih file mana yang jadi diproses.
    """
    if count <= MAX_FILES_PER_BATCH:
        return None
    return (
        f"Kamu memilih {count} file sekaligus, sedangkan batasnya "
        f"{MAX_FILES_PER_BATCH} file per unggah. Tidak ada yang diproses. "
        f"Silakan pilih ulang maksimal {MAX_FILES_PER_BATCH} file."
    )


def _short_reason(exc: Exception) -> str:
    """Pesan error yang sudah ramah dan cukup pendek untuk sidebar."""
    reason = errors.friendly_message(exc)
    return reason if len(reason) <= 110 else reason[:107] + "…"


def delete_document(document_id: int) -> None:
    """Hapus dokumen: chunk di Qdrant dulu, baru metadata dan teksnya di SQLite."""
    rag.remove_document(document_id)
    db.delete_document(document_id)


def _clean_text(raw: str) -> str:
    """Rapikan hasil ekstraksi: spasi berlebih dan baris kosong bertumpuk."""
    text = raw.replace("\r\n", "\n").replace("\r", "\n")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r" *\n *", "\n", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# --------------------------------------------------------------------------- #
# Tahap 2 — validasi
# --------------------------------------------------------------------------- #
def validate(filename: str, size_on_disk: int) -> str:
    """Pastikan tipe dan ukuran file memenuhi syarat.

    Mengembalikan ekstensi file kalau lolos, atau melempar ``RejectedFile``
    dengan pesan yang siap ditampilkan ke pengguna.
    """
    extension = file_extension(filename)

    if not extension:
        raise RejectedFile("File tidak punya ekstensi, jadi tipenya tidak bisa dikenali.")

    if extension not in ALLOWED_TYPES:
        raise RejectedFile(f"Tipe .{extension} belum didukung.")

    if size_on_disk == 0:
        raise RejectedFile("File-nya kosong (0 byte), tidak ada yang bisa diproses.")

    if size_on_disk > MAX_FILE_SIZE:
        raise RejectedFile(
            f"Ukurannya {human_size(size_on_disk)}, melebihi batas {MAX_FILE_SIZE_MB} MB per file."
        )

    return extension


# --------------------------------------------------------------------------- #
# Tahap 3 — parsing per tipe
# --------------------------------------------------------------------------- #
def _parse_pdf(path: Path) -> str:
    from pypdf import PdfReader
    from pypdf.errors import PdfReadError

    try:
        reader = PdfReader(path)
        if reader.is_encrypted:
            # Beberapa PDF terenkripsi dengan password kosong dan masih bisa dibuka.
            try:
                reader.decrypt("")
            except Exception as exc:
                raise ParseFailed("PDF-nya terkunci dengan password.") from exc
        pages = [(page.extract_text() or "") for page in reader.pages]
    except ParseFailed:
        raise
    except (PdfReadError, Exception) as exc:
        raise ParseFailed("PDF-nya tidak bisa dibaca, kemungkinan rusak atau terkunci.") from exc

    return "\n\n".join(page for page in pages if page.strip())


def _parse_docx(path: Path) -> str:
    import docx

    try:
        document = docx.Document(str(path))
    except Exception as exc:
        raise ParseFailed("File DOCX-nya tidak bisa dibuka, kemungkinan rusak.") from exc

    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            if any(cells):
                parts.append(" | ".join(cells))
    return "\n".join(parts)


def _parse_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    try:
        workbook = load_workbook(str(path), read_only=True, data_only=True)
    except Exception as exc:
        raise ParseFailed("File XLSX-nya tidak bisa dibuka, kemungkinan rusak.") from exc

    parts: list[str] = []
    try:
        for sheet in workbook.worksheets:
            parts.append(f"# Sheet: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(value) for value in row if value is not None]
                if cells:
                    parts.append(" | ".join(cells))
    finally:
        workbook.close()
    return "\n".join(parts)


def _parse_pptx(path: Path) -> str:
    from pptx import Presentation

    try:
        presentation = Presentation(str(path))
    except Exception as exc:
        raise ParseFailed("File PPTX-nya tidak bisa dibuka, kemungkinan rusak.") from exc

    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"# Slide {index}")
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            parts.append(f"[Catatan] {slide.notes_slide.notes_text_frame.text}")
    return "\n".join(parts)


PARSERS: dict[str, Callable[[Path], str]] = {
    "pdf": _parse_pdf,
    "docx": _parse_docx,
    "xlsx": _parse_xlsx,
    "pptx": _parse_pptx,
}


def parse(path: Path, extension: str) -> str:
    """Ekstrak teks dari file dokumen di ``path``."""
    parser = PARSERS[extension]
    return _clean_text(parser(path))


def describe(image_bytes: bytes, filename: str) -> str:
    """Baca isi gambar lewat model vision, hasilnya jadi teks dokumen."""
    return _clean_text(vision.describe_image(image_bytes, filename))


# --------------------------------------------------------------------------- #
# Pipeline lengkap
# --------------------------------------------------------------------------- #
def ingest(
    filename: str,
    data: bytes,
    *,
    progress: Callable[[str], None] | None = None,
) -> IngestResult:
    """Jalankan lima tahap pemrosesan untuk satu file.

    Selalu mengembalikan ``IngestResult`` — tidak pernah melempar exception
    mentah ke pemanggil — dan file temp-nya dijamin terhapus lewat ``finally``.
    """

    def step(message: str) -> None:
        if progress is not None:
            progress(message)

    temp_path: Path | None = None
    try:
        # Tahap 1 — simpan sementara ke folder temp
        step("Menyimpan file ke folder temp…")
        suffix = Path(filename).suffix
        with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as handle:
            handle.write(data)
            temp_path = Path(handle.name)

        # Tahap 2 — validasi tipe dan ukuran
        step("Memvalidasi tipe dan ukuran…")
        size_on_disk = temp_path.stat().st_size
        extension = validate(filename, size_on_disk)

        # Tahap 3 — parsing jadi teks. Untuk gambar, "teks"-nya adalah
        # deskripsi dari model vision Groq, termasuk tulisan yang terbaca
        # di dalam gambar. Setelah ini alurnya sama persis dengan dokumen.
        vision_note = ""
        if extension in IMAGE_TYPES:
            step(f"Membaca gambar dengan {get_settings().vision_model}…")
            try:
                text = describe(data, filename)
                status = STATUS_PROCESSED if text else STATUS_NO_TEXT
            except vision.VisionError as exc:
                text = ""
                status = STATUS_PENDING_EXTRACTION
                vision_note = str(exc)
        else:
            step("Mengekstrak teks…")
            text = parse(temp_path, extension)
            status = STATUS_PROCESSED if text else STATUS_NO_TEXT

        # Tahap 4 — simpan metadata dan teks ke SQLite
        step("Menyimpan metadata dan teks ke SQLite…")
        document_id = db.save_document(
            filename=filename,
            filetype=extension,
            filesize=size_on_disk,
            status=status,
            text=text,
        )

        # Tahap 5 — file asli dihapus dari disk, sebelum pekerjaan embedding
        # yang bisa memakan waktu. Blok finally di bawah tetap jadi jaring
        # pengaman kalau tahap sebelumnya sempat gagal.
        step("Menghapus file asli dari disk…")
        temp_path.unlink(missing_ok=True)
        temp_path = None

        # Tahap 6 — chunking, embedding, lalu simpan vector-nya ke Qdrant
        chunk_count = 0
        index_note = ""
        if text:
            try:
                chunk_count = rag.index_document(
                    document_id, filename, text, progress=step
                )
                if chunk_count:
                    status = STATUS_INDEXED
                    db.set_document_status(document_id, status)
            except Exception as exc:  # noqa: BLE001 - indeks gagal, dokumen tetap tersimpan
                index_note = f" Belum terindeks untuk RAG: {_short_reason(exc)}"

        if status == STATUS_PENDING_EXTRACTION:
            message = f"tersimpan sebagai metadata saja. {vision_note}".strip()
        elif status == STATUS_NO_TEXT:
            message = "tersimpan, tapi tidak ada teks di dalamnya."
        elif status == STATUS_INDEXED:
            source = "deskripsi gambar" if extension in IMAGE_TYPES else "teks"
            message = (
                f"terindeks, {thousands(len(text))} karakter {source} "
                f"jadi {chunk_count} chunk."
            )
        else:
            message = f"diproses, {thousands(len(text))} karakter tersimpan.{index_note}"

        return IngestResult(
            filename=filename,
            ok=True,
            message=message,
            status=status,
            char_count=len(text),
            document_id=document_id,
            chunk_count=chunk_count,
        )

    except RejectedFile as exc:
        return IngestResult(filename=filename, ok=False, message=str(exc), status=STATUS_FAILED)
    except ParseFailed as exc:
        return IngestResult(filename=filename, ok=False, message=str(exc), status=STATUS_FAILED)
    except Exception:  # noqa: BLE001 - error tak terduga tetap disampaikan dengan sopan
        return IngestResult(
            filename=filename,
            ok=False,
            message="Ada kendala tak terduga saat memprosesnya. Silakan coba lagi.",
            status=STATUS_FAILED,
        )
    finally:
        # Jaring pengaman: kalau pipeline berhenti sebelum tahap 5, file
        # sementaranya tetap dihapus.
        if temp_path is not None:
            step("Menghapus file asli dari disk…")
            temp_path.unlink(missing_ok=True)
